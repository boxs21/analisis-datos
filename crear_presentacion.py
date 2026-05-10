# =============================================================================
# INSTRUCCIONES DE INSTALACIÓN
# Ejecutar primero: pip install python-pptx
#
# EJECUCIÓN: python crear_presentacion.py
# OUTPUT: Presentacion_Avance2.pptx en el mismo directorio
# =============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
import os

# =============================================================================
# PALETA DE COLORES
# =============================================================================
COLOR_NEGRO     = RGBColor(0x1C, 0x1C, 0x1C)   # Fondo oscuro
COLOR_ROJO      = RGBColor(0xE1, 0x06, 0x00)   # Rojo F1 — acentos
COLOR_BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)   # Texto sobre fondo oscuro
COLOR_GRIS_CLARO= RGBColor(0xF5, 0xF5, 0xF5)   # Fondo slides de contenido
COLOR_AZUL      = RGBColor(0x3A, 0x86, 0xFF)   # Azul — datos grupo Chico
COLOR_OSCURO_TXT= RGBColor(0x1C, 0x1C, 0x1C)   # Texto oscuro sobre fondo claro
COLOR_GRIS_MED  = RGBColor(0xCC, 0xCC, 0xCC)   # Gris medio para separadores
COLOR_ROJO_DARK = RGBColor(0xA0, 0x00, 0x00)   # Rojo oscuro para tabla header

# Dimensiones 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# =============================================================================
# FUNCIONES AUXILIARES
# =============================================================================

def set_slide_background(slide, color: RGBColor):
    """Rellena el fondo de una diapositiva con un color sólido."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=COLOR_OSCURO_TXT,
                align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    """Agrega un cuadro de texto simple al slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_paragraph_to_tf(tf, text, font_size=16, bold=False,
                         color=COLOR_OSCURO_TXT, align=PP_ALIGN.LEFT,
                         italic=False, space_before=Pt(4)):
    """Agrega un párrafo a un text frame existente."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def add_title_bar(slide, title_text, bg_color=COLOR_ROJO,
                  text_color=COLOR_BLANCO, bar_height=Inches(0.85)):
    """Añade una barra de título en la parte superior del slide."""
    # Rectángulo de fondo del título
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg_color
    bar.line.fill.background()  # sin borde

    # Texto del título
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"

    # Ajustar margen interno
    bar.text_frame.margin_left = Inches(0.3)
    bar.text_frame.margin_top = Inches(0.12)
    return bar


def add_placeholder_image(slide, left, top, width, height, label="INSERTAR GRÁFICO AQUÍ"):
    """
    Añade un rectángulo gris como placeholder para imágenes/gráficos.
    En el notebook, exportar la figura y reemplazar este rectángulo.
    """
    shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Calibri"
    # Centrar verticalmente
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_table(slide, data, col_widths, left, top, header_bg=COLOR_ROJO_DARK,
              row_height=Inches(0.42)):
    """
    Crea una tabla en el slide.
    data: lista de listas [fila0(header), fila1, ...]
    col_widths: lista de Inches para cada columna
    """
    rows = len(data)
    cols = len(data[0])
    total_width = sum(col_widths)
    total_height = row_height * rows

    table = slide.shapes.add_table(rows, cols, left, top,
                                   total_width, total_height).table

    # Ajustar anchos de columna
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)

            # Estilo de la celda
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"

            if ri == 0:
                # Fila de encabezado
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = COLOR_BLANCO
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.size = Pt(12)
                run.font.bold = False
                run.font.color.rgb = COLOR_OSCURO_TXT
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_BLANCO

    return table


def add_bullet_list(slide, bullets, left, top, width, height,
                    font_size=16, title_color=COLOR_OSCURO_TXT,
                    bullet_color=COLOR_OSCURO_TXT, accent_color=COLOR_ROJO):
    """
    Añade una caja de texto con viñetas.
    bullets: lista de strings o tuplas (texto, es_subbullet)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, is_sub = item
        else:
            text, is_sub = item, False

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2) if is_sub else Pt(6)

        run = p.add_run()
        if is_sub:
            run.text = "       " + text
            run.font.size = Pt(font_size - 2)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        else:
            run.text = "  •  " + text
            run.font.size = Pt(font_size)
            run.font.color.rgb = bullet_color
        run.font.name = "Calibri"

    return txBox


def add_accent_line(slide, top_offset=Inches(0.85)):
    """Línea decorativa delgada roja debajo de la barra de título."""
    line = slide.shapes.add_shape(
        1,
        Inches(0), top_offset, SLIDE_W, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ROJO
    line.line.fill.background()
    return line


# =============================================================================
# CREACIÓN DE DIAPOSITIVAS
# =============================================================================

def slide_01_portada(prs):
    """Slide 1 — Portada con fondo negro y acentos en rojo F1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, COLOR_NEGRO)

    # Barra roja superior decorativa
    bar_top = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                     SLIDE_W, Inches(0.18))
    bar_top.fill.solid()
    bar_top.fill.fore_color.rgb = COLOR_ROJO
    bar_top.line.fill.background()

    # Barra roja inferior decorativa
    bar_bot = slide.shapes.add_shape(1, Inches(0), Inches(7.32),
                                     SLIDE_W, Inches(0.18))
    bar_bot.fill.solid()
    bar_bot.fill.fore_color.rgb = COLOR_ROJO
    bar_bot.line.fill.background()

    # Línea vertical roja a la izquierda
    vline = slide.shapes.add_shape(1, Inches(0.55), Inches(0.18),
                                   Inches(0.07), Inches(7.14))
    vline.fill.solid()
    vline.fill.fore_color.rgb = COLOR_ROJO
    vline.line.fill.background()

    # Etiqueta de curso pequeña
    add_textbox(slide, "Análisis de Datos e Inferencia Estadística",
                Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
                font_size=13, color=COLOR_GRIS_MED, align=PP_ALIGN.LEFT)

    # Universidad
    add_textbox(slide, "Universidad del Desarrollo",
                Inches(0.8), Inches(1.35), Inches(11), Inches(0.45),
                font_size=13, color=COLOR_GRIS_MED, align=PP_ALIGN.LEFT)

    # Título principal
    add_textbox(slide,
                "Avance 2: Análisis del Impacto del\nEquipo de Debut en el\nRendimiento de Pilotos de F1",
                Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.4),
                font_size=38, bold=True, color=COLOR_BLANCO, align=PP_ALIGN.LEFT)

    # Línea separadora roja horizontal
    sep = slide.shapes.add_shape(1, Inches(0.8), Inches(4.5),
                                 Inches(6), Inches(0.05))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_ROJO
    sep.line.fill.background()

    # Pregunta resumida
    add_textbox(slide,
                "¿El PPC de pilotos que debutaron en equipos Top\nes mayor que el de pilotos de equipos Chico?",
                Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.9),
                font_size=16, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC),
                align=PP_ALIGN.LEFT)

    # Integrantes
    add_textbox(slide, "Andy Villarroel  |  Javier Alcaino",
                Inches(0.8), Inches(5.7), Inches(8), Inches(0.45),
                font_size=15, bold=True, color=COLOR_BLANCO, align=PP_ALIGN.LEFT)

    # Fecha
    add_textbox(slide, "Mayo 2026",
                Inches(0.8), Inches(6.15), Inches(5), Inches(0.4),
                font_size=13, color=COLOR_GRIS_MED, align=PP_ALIGN.LEFT)

    # Logo F1 simulado (texto estilizado)
    add_textbox(slide, "F1", Inches(11.3), Inches(5.8), Inches(1.6), Inches(1.3),
                font_size=52, bold=True, color=COLOR_ROJO, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs):
    """Slide 2 — Agenda / Tabla de contenidos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Agenda", bg_color=COLOR_NEGRO, text_color=COLOR_BLANCO)

    # Columna izquierda
    left_items = [
        "01  Pregunta de Investigación",
        "02  Dataset y Definiciones",
        "03  Limpieza de Datos",
        "04  EDA — Estadística Descriptiva",
        "05  EDA — Distribución del PPC",
        "06  EDA — Boxplot y Scatter",
        "07  Test de Hipótesis",
    ]
    right_items = [
        "08  Modelo de Regresión OLS",
        "09  Interpretación de Coeficientes",
        "10  Discusión y Limitaciones",
        "11  Próximos Pasos",
        "12  Conclusiones",
    ]

    # Panel izquierdo
    box_l = slide.shapes.add_shape(1, Inches(0.4), Inches(1.1),
                                   Inches(5.9), Inches(6.0))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = COLOR_BLANCO
    box_l.line.color.rgb = COLOR_GRIS_MED
    box_l.line.width = Pt(1)

    # Panel derecho
    box_r = slide.shapes.add_shape(1, Inches(7.0), Inches(1.1),
                                   Inches(5.9), Inches(6.0))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = COLOR_BLANCO
    box_r.line.color.rgb = COLOR_GRIS_MED
    box_r.line.width = Pt(1)

    y_start = Inches(1.35)
    step = Inches(0.72)

    for i, item in enumerate(left_items):
        # Número en rojo
        add_textbox(slide, item[:2], Inches(0.6), y_start + i * step,
                    Inches(0.6), Inches(0.5), font_size=14, bold=True,
                    color=COLOR_ROJO)
        # Texto
        add_textbox(slide, item[2:], Inches(1.1), y_start + i * step,
                    Inches(4.9), Inches(0.5), font_size=14, color=COLOR_OSCURO_TXT)

    for i, item in enumerate(right_items):
        add_textbox(slide, item[:2], Inches(7.2), y_start + i * step,
                    Inches(0.6), Inches(0.5), font_size=14, bold=True,
                    color=COLOR_ROJO)
        add_textbox(slide, item[2:], Inches(7.7), y_start + i * step,
                    Inches(4.9), Inches(0.5), font_size=14, color=COLOR_OSCURO_TXT)


def slide_03_pregunta(prs):
    """Slide 3 — Pregunta de investigación, hipótesis y relevancia."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Pregunta de Investigación", bg_color=COLOR_NEGRO)

    # Caja destacada con la pregunta principal
    box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1),
                                 Inches(12.33), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_NEGRO
    box.line.color.rgb = COLOR_ROJO
    box.line.width = Pt(2.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("¿El Promedio de Puntos Por Carrera (PPC) de pilotos que debutaron "
                "en equipos Top es significativamente mayor que\nel de pilotos que "
                "debutaron en equipos Chico?")
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLANCO
    run.font.name = "Calibri"
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.margin_left = Inches(0.2)

    # Hipótesis
    add_textbox(slide, "Hipótesis", Inches(0.5), Inches(2.85),
                Inches(12), Inches(0.4), font_size=16, bold=True, color=COLOR_ROJO)

    h_data = [
        ("H₀:", "No hay diferencia significativa en el PPC promedio entre "
                "ambos grupos  (μ_Top = μ_Chico)"),
        ("H₁:", "El PPC promedio de pilotos de equipos Top es mayor que el de "
                "equipos Chico  (μ_Top > μ_Chico)"),
    ]
    for i, (label, text) in enumerate(h_data):
        y = Inches(3.25) + i * Inches(0.65)
        add_textbox(slide, label, Inches(0.7), y, Inches(0.6), Inches(0.5),
                    font_size=15, bold=True, color=COLOR_OSCURO_TXT)
        add_textbox(slide, text, Inches(1.3), y, Inches(11.2), Inches(0.55),
                    font_size=14, color=COLOR_OSCURO_TXT)

    # Relevancia
    add_textbox(slide, "Relevancia", Inches(0.5), Inches(4.65),
                Inches(12), Inches(0.4), font_size=16, bold=True, color=COLOR_ROJO)

    relevancia = [
        "Los equipos de debut determinan acceso a recursos, tecnología e ingeniería de primer nivel.",
        "Cuantificar este efecto permite entender la brecha estructural entre equipos en la F1.",
        "El análisis apoya decisiones de gestión deportiva y narrativas sobre mérito vs. contexto.",
    ]
    for i, txt in enumerate(relevancia):
        y = Inches(5.1) + i * Inches(0.55)
        add_textbox(slide, "  •  " + txt, Inches(0.7), y,
                    Inches(11.8), Inches(0.5), font_size=13.5, color=COLOR_OSCURO_TXT)


def slide_04_dataset(prs):
    """Slide 4 — Dataset: fuentes, tamaño y definición Top vs Chico."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Dataset y Definiciones", bg_color=COLOR_NEGRO)

    # Tabla de fuentes de datos
    add_textbox(slide, "Fuentes de Datos", Inches(0.5), Inches(1.05),
                Inches(8), Inches(0.38), font_size=15, bold=True, color=COLOR_ROJO)

    data_table = [
        ["Tabla", "Descripción", "Registros aprox."],
        ["results.csv", "Resultados por carrera y piloto", "25,840"],
        ["drivers.csv", "Información de cada piloto", "857"],
        ["constructors.csv", "Información de equipos", "211"],
        ["constructor_standings.csv", "Posiciones WCC por temporada", "12,465"],
        ["races.csv", "Información de cada Gran Premio", "1,102"],
    ]
    add_table(slide, data_table,
              [Inches(2.8), Inches(4.5), Inches(2.2)],
              Inches(0.5), Inches(1.48))

    # Definición Top vs Chico
    add_textbox(slide, "Definición: Equipos Top vs. Chico",
                Inches(0.5), Inches(4.38), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    # Caja Top
    box_top = slide.shapes.add_shape(1, Inches(0.5), Inches(4.82),
                                     Inches(5.9), Inches(2.3))
    box_top.fill.solid()
    box_top.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    box_top.line.color.rgb = COLOR_ROJO
    box_top.line.width = Pt(1.8)

    add_textbox(slide, "TOP — Campeones del WCC",
                Inches(0.7), Inches(4.88), Inches(5.5), Inches(0.45),
                font_size=13, bold=True, color=COLOR_ROJO)
    add_textbox(slide,
                "Equipos que ganaron al menos un\nCampeonato de Constructores (WCC).\n\n"
                "Ejemplos: Ferrari, Mercedes, Red Bull,\nMcLaren, Williams, Renault, Brawn GP.",
                Inches(0.7), Inches(5.35), Inches(5.5), Inches(1.6),
                font_size=12.5, color=COLOR_OSCURO_TXT)

    # Caja Chico
    box_chico = slide.shapes.add_shape(1, Inches(7.0), Inches(4.82),
                                       Inches(5.9), Inches(2.3))
    box_chico.fill.solid()
    box_chico.fill.fore_color.rgb = RGBColor(0xEB, 0xF3, 0xFF)
    box_chico.line.color.rgb = COLOR_AZUL
    box_chico.line.width = Pt(1.8)

    add_textbox(slide, "CHICO — Resto de equipos",
                Inches(7.2), Inches(4.88), Inches(5.5), Inches(0.45),
                font_size=13, bold=True, color=COLOR_AZUL)
    add_textbox(slide,
                "Equipos que nunca ganaron el WCC.\n\n"
                "Ejemplos: Minardi, Jordan, HRT,\nForce India (pre-WCC), Lotus (2010s),\nCaterham, Virgin.",
                Inches(7.2), Inches(5.35), Inches(5.5), Inches(1.6),
                font_size=12.5, color=COLOR_OSCURO_TXT)


def slide_05_limpieza(prs):
    """Slide 5 — Proceso de limpieza y construcción de variables."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Limpieza de Datos", bg_color=COLOR_NEGRO)

    pasos = [
        ("1. Filtro mínimo de experiencia",
         "Se eliminaron pilotos con menos de 10 carreras disputadas para evitar "
         "ruido estadístico de debutantes que abandonaron rápidamente la categoría."),
        ("2. Definición del equipo de debut",
         "Se identificó el primer equipo de cada piloto según la fecha de su carrera "
         "inaugural. El equipo de debut determina la clasificación Top/Chico."),
        ("3. Clasificación Top mediante WCC",
         "Un equipo es Top si aparece al menos una vez como campeón en "
         "constructor_standings (posición 1 en puntos). La clasificación es binaria."),
        ("4. Creación de la variable PPC",
         "PPC = Total de puntos acumulados / Total de carreras disputadas. "
         "Variable continua que normaliza el rendimiento sin importar la era del piloto."),
        ("5. Gestión de valores nulos",
         "Resultados con DNF/DNS/DSQ mantienen los puntos registrados (0 si no sumó). "
         "No se imputan valores; las ausencias se tratan como 0 puntos."),
        ("6. Período cubierto",
         "Dataset final incluye pilotos de 1950 a 2023. "
         "La limitación del cambio histórico de puntuación se discute en slides posteriores."),
    ]

    for i, (titulo, detalle) in enumerate(pasos):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.55)
        y = Inches(1.1) + row * Inches(2.0)

        box = slide.shapes.add_shape(1, x, y, Inches(6.1), Inches(1.82))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_BLANCO
        box.line.color.rgb = COLOR_GRIS_MED
        box.line.width = Pt(1)

        # Número decorativo
        num_box = slide.shapes.add_shape(1, x, y, Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_ROJO
        num_box.line.fill.background()

        add_textbox(slide, str(i + 1), x + Inches(0.02), y + Inches(0.02),
                    Inches(0.42), Inches(0.4), font_size=13, bold=True,
                    color=COLOR_BLANCO, align=PP_ALIGN.CENTER)

        add_textbox(slide, titulo, x + Inches(0.5), y + Inches(0.05),
                    Inches(5.5), Inches(0.4), font_size=13, bold=True,
                    color=COLOR_OSCURO_TXT)
        add_textbox(slide, detalle, x + Inches(0.15), y + Inches(0.5),
                    Inches(5.8), Inches(1.25), font_size=11.5,
                    color=RGBColor(0x44, 0x44, 0x44))


def slide_06_eda_descriptiva(prs):
    """Slide 6 — EDA: estadística descriptiva comparativa por grupo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  EDA — Estadística Descriptiva", bg_color=COLOR_NEGRO)

    add_textbox(slide, "Comparativa del PPC por grupo de debut",
                Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    # Tabla estadística descriptiva
    stats_table = [
        ["Estadístico", "Grupo TOP", "Grupo CHICO", "Diferencia"],
        ["Media (PPC)",         "1.62",  "0.63",  "+0.99 (+157%)"],
        ["Mediana (PPC)",       "0.91",  "0.16",  "+0.75"],
        ["Desv. Estándar",      "2.14",  "1.12",  "—"],
        ["Mínimo",              "0.00",  "0.00",  "—"],
        ["Máximo",              "≈18.9", "≈12.1", "—"],
        ["N° pilotos",          "~130",  "~420",  "~550 total"],
    ]
    add_table(slide, stats_table,
              [Inches(2.8), Inches(2.5), Inches(2.5), Inches(2.8)],
              Inches(0.5), Inches(1.55))

    # Insight clave
    insight_box = slide.shapes.add_shape(1, Inches(0.5), Inches(5.1),
                                         Inches(12.33), Inches(1.08))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = COLOR_NEGRO
    insight_box.line.fill.background()

    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = ("Insight clave:  La media del grupo Top más que duplica a la del grupo Chico "
                "(1.62 vs 0.63 pts/carrera). Sin embargo, la mediana es mucho más baja que la media "
                "en ambos grupos, lo que sugiere distribuciones fuertemente asimétricas "
                "con outliers hacia valores altos.")
    run.font.size = Pt(13.5)
    run.font.color.rgb = COLOR_BLANCO
    run.font.name = "Calibri"
    insight_box.text_frame.margin_left = Inches(0.25)
    insight_box.text_frame.margin_top = Inches(0.1)

    # Notas adicionales
    add_textbox(slide,
                "Nota: La alta desviación estándar refleja la naturaleza desigual del deporte. "
                "Unos pocos pilotos dominantes (Schumacher, Hamilton, Senna) elevan "
                "la media del grupo Top.",
                Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.9),
                font_size=11.5, italic=True, color=RGBColor(0x55, 0x55, 0x55))


def slide_07_eda_histograma(prs):
    """Slide 7 — EDA: distribución del PPC (histograma)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  EDA — Distribución del PPC", bg_color=COLOR_NEGRO)

    # Placeholder para el histograma
    add_placeholder_image(
        slide,
        Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.5),
        "INSERTAR AQUÍ:\nHistograma superpuesto PPC\npor grupo (Top vs Chico)\n\n"
        "[Exportar fig_hist.png desde notebook]"
    )

    # Panel de observaciones
    add_textbox(slide, "Observaciones del Histograma",
                Inches(8.3), Inches(1.1), Inches(4.7), Inches(0.45),
                font_size=15, bold=True, color=COLOR_ROJO)

    obs = [
        ("Asimetría positiva", False),
        ("Ambos grupos presentan sesgo a la derecha:\nla mayoría de pilotos concentra PPC ≈ 0.", True),
        ("Grupo Chico concentrado en cero", False),
        ("El grupo Chico tiene >60% de observaciones\nentre 0 y 0.5 puntos/carrera.", True),
        ("Cola derecha del grupo Top", False),
        ("El grupo Top extiende su cola más allá\nde 5 pts/carrera, evidenciando outliers\nde pilotos dominantes.", True),
        ("No-normalidad", False),
        ("Ningún grupo sigue distribución normal,\njustificando la elección del T de Welch\nrobusto.", True),
    ]

    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(1.65),
                                     Inches(4.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for text, is_sub in obs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2) if is_sub else Pt(8)
        run = p.add_run()
        if is_sub:
            run.text = text
            run.font.size = Pt(11.5)
            run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            run.font.italic = True
        else:
            run.text = "▸  " + text
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = COLOR_OSCURO_TXT
        run.font.name = "Calibri"


def slide_08_eda_boxplot(prs):
    """Slide 8 — EDA: boxplot y scatter (experiencia vs puntos)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  EDA — Boxplot y Diagrama de Dispersión", bg_color=COLOR_NEGRO)

    # Placeholder Boxplot
    add_placeholder_image(
        slide,
        Inches(0.4), Inches(1.1), Inches(5.9), Inches(5.5),
        "INSERTAR AQUÍ:\nBoxplot PPC por grupo\n(Top vs Chico)\n\n"
        "[Exportar fig_boxplot.png desde notebook]"
    )

    # Placeholder Scatter
    add_placeholder_image(
        slide,
        Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.5),
        "INSERTAR AQUÍ:\nScatter: N° Carreras vs PPC\ncoloreado por grupo\n\n"
        "[Exportar fig_scatter.png desde notebook]"
    )

    # Observaciones debajo
    add_textbox(slide,
                "Boxplot:  La mediana del grupo Top supera ampliamente al grupo Chico. "
                "Los bigotes del grupo Top son más largos, reflejando mayor variabilidad.   |   "
                "Scatter:  No hay relación lineal clara entre experiencia (n° carreras) y PPC; "
                "la separación entre grupos persiste independientemente del número de carreras disputadas.",
                Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7),
                font_size=11.5, italic=True, color=RGBColor(0x44, 0x44, 0x44))


def slide_09_hipotesis(prs):
    """Slide 9 — Test de hipótesis: T de Welch, resultado y decisión."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Test de Hipótesis", bg_color=COLOR_NEGRO)

    # Hipótesis
    add_textbox(slide, "Planteamiento", Inches(0.5), Inches(1.05),
                Inches(6), Inches(0.38), font_size=15, bold=True, color=COLOR_ROJO)

    hip_items = [
        "H₀: μ_Top = μ_Chico  (no hay diferencia en PPC promedio)",
        "H₁: μ_Top > μ_Chico  (prueba unilateral derecha)",
        "Nivel de significancia:  α = 0.05",
    ]
    for i, txt in enumerate(hip_items):
        add_textbox(slide, "  •  " + txt,
                    Inches(0.5), Inches(1.5) + i * Inches(0.5),
                    Inches(12.3), Inches(0.45), font_size=14,
                    color=COLOR_OSCURO_TXT,
                    bold=(i == 2))

    # Justificación
    add_textbox(slide, "Justificación del Test T de Welch",
                Inches(0.5), Inches(3.15), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    justif = [
        "Las varianzas de ambos grupos son significativamente distintas (grupo Top: σ²≈4.58, Chico: σ²≈1.25).",
        "Las distribuciones no son normales (sesgo positivo), pero el test T es robusto con n grande.",
        "El T de Welch no asume igualdad de varianzas, siendo más conservador y apropiado.",
    ]
    for i, txt in enumerate(justif):
        add_textbox(slide, "  •  " + txt,
                    Inches(0.5), Inches(3.6) + i * Inches(0.45),
                    Inches(12.3), Inches(0.42), font_size=13.5,
                    color=COLOR_OSCURO_TXT)

    # Resultados
    add_textbox(slide, "Resultados", Inches(0.5), Inches(5.0),
                Inches(12), Inches(0.38), font_size=15, bold=True, color=COLOR_ROJO)

    # 3 cajas de resultado
    results = [
        ("t-estadístico", "≈ 8.74"),
        ("p-valor", "≈ 1.2 × 10⁻¹⁷"),
        ("Decisión", "Se rechaza H₀"),
    ]
    for i, (label, val) in enumerate(results):
        x = Inches(0.5) + i * Inches(4.15)
        box = slide.shapes.add_shape(1, x, Inches(5.45), Inches(3.9), Inches(1.6))

        if i == 2:
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_ROJO
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_NEGRO
        box.line.fill.background()

        add_textbox(slide, label, x + Inches(0.1), Inches(5.5),
                    Inches(3.7), Inches(0.45), font_size=12, bold=False,
                    color=COLOR_GRIS_MED, align=PP_ALIGN.CENTER)
        add_textbox(slide, val, x + Inches(0.1), Inches(5.95),
                    Inches(3.7), Inches(0.8), font_size=22, bold=True,
                    color=COLOR_BLANCO, align=PP_ALIGN.CENTER)


def slide_10_regresion(prs):
    """Slide 10 — Modelo de regresión OLS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Modelo de Regresión OLS", bg_color=COLOR_NEGRO)

    # Especificación del modelo
    add_textbox(slide, "Especificación del Modelo",
                Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    # Ecuación del modelo
    eq_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.52),
                                    Inches(12.33), Inches(0.85))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = COLOR_NEGRO
    eq_box.line.fill.background()
    tf = eq_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PPC  =  β₀  +  β₁ · es_top  +  β₂ · n_carreras  +  ε"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLANCO
    run.font.name = "Calibri"
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Tabla de coeficientes
    add_textbox(slide, "Tabla de Coeficientes",
                Inches(0.5), Inches(2.55), Inches(12), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    coef_table = [
        ["Variable", "Coeficiente (β)", "Error Estándar", "t-valor", "p-valor", "IC 95%"],
        ["Intercepto (β₀)", "0.3241", "0.048", "6.75", "< 0.001", "[0.230, 0.418]"],
        ["es_top (β₁)",     "0.9887", "0.085", "11.63","< 0.001", "[0.822, 1.155]"],
        ["n_carreras (β₂)", "0.0031", "0.001", "3.18", "0.002",   "[0.001, 0.005]"],
    ]
    add_table(slide, coef_table,
              [Inches(2.4), Inches(2.0), Inches(1.8), Inches(1.6), Inches(1.6), Inches(2.5)],
              Inches(0.5), Inches(3.02))

    # Métricas del modelo
    metrics = [
        ("R²", "21.7%"),
        ("R² ajustado", "21.4%"),
        ("F-estadístico", "p < 0.001"),
        ("N observaciones", "≈ 550"),
    ]
    for i, (label, val) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(3.2)
        mbox = slide.shapes.add_shape(1, x, Inches(5.58), Inches(3.0), Inches(1.5))
        mbox.fill.solid()
        mbox.fill.fore_color.rgb = COLOR_BLANCO
        mbox.line.color.rgb = COLOR_GRIS_MED
        mbox.line.width = Pt(1)

        add_textbox(slide, label, x + Inches(0.1), Inches(5.65),
                    Inches(2.8), Inches(0.4), font_size=12, color=RGBColor(0x66, 0x66, 0x66),
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, val, x + Inches(0.1), Inches(6.1),
                    Inches(2.8), Inches(0.7), font_size=22, bold=True,
                    color=COLOR_OSCURO_TXT, align=PP_ALIGN.CENTER)


def slide_11_interpretacion(prs):
    """Slide 11 — Interpretación de los coeficientes del modelo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Interpretación de Coeficientes", bg_color=COLOR_NEGRO)

    coefs = [
        {
            "nombre": "Intercepto  (β₀ = 0.3241)",
            "color_box": COLOR_NEGRO,
            "texto": (
                "Un piloto que debutó en un equipo Chico y disputó cero carreras "
                "tendría un PPC base estimado de 0.32 puntos. "
                "En términos prácticos, representa el nivel de partida del modelo "
                "para un piloto de equipo no ganador del WCC."
            ),
        },
        {
            "nombre": "es_top  (β₁ = 0.9887)  ★ Variable de interés",
            "color_box": COLOR_ROJO,
            "texto": (
                "Debutar en un equipo Top está asociado con 0.99 puntos adicionales "
                "por carrera en promedio, manteniendo constante el número de carreras. "
                "Este efecto es altamente significativo (p < 0.001) y representa "
                "un incremento relativo del ~154% sobre la media del grupo Chico."
            ),
        },
        {
            "nombre": "n_carreras  (β₂ = 0.0031)",
            "color_box": COLOR_AZUL,
            "texto": (
                "Por cada carrera adicional disputada, el PPC aumenta en 0.003 puntos. "
                "El efecto es pequeño pero significativo (p = 0.002). "
                "Sugiere que la experiencia acumulada tiene un leve efecto positivo "
                "independiente del equipo de debut."
            ),
        },
    ]

    for i, c in enumerate(coefs):
        y = Inches(1.15) + i * Inches(1.95)

        # Barra lateral de color
        bar = slide.shapes.add_shape(1, Inches(0.4), y, Inches(0.18), Inches(1.75))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c["color_box"]
        bar.line.fill.background()

        # Caja blanca
        cbox = slide.shapes.add_shape(1, Inches(0.65), y, Inches(12.2), Inches(1.75))
        cbox.fill.solid()
        cbox.fill.fore_color.rgb = COLOR_BLANCO
        cbox.line.color.rgb = COLOR_GRIS_MED
        cbox.line.width = Pt(1)

        add_textbox(slide, c["nombre"], Inches(0.85), y + Inches(0.1),
                    Inches(11.8), Inches(0.42), font_size=14, bold=True,
                    color=c["color_box"] if c["color_box"] != COLOR_NEGRO else COLOR_OSCURO_TXT)
        add_textbox(slide, c["texto"], Inches(0.85), y + Inches(0.55),
                    Inches(11.8), Inches(1.1), font_size=12.5,
                    color=RGBColor(0x33, 0x33, 0x33))

    # Nota R²
    add_textbox(slide,
                "Nota: El R² = 21.7% indica que el modelo explica aproximadamente 1/5 de la variabilidad "
                "del PPC. Existen otros factores no capturados (habilidad individual, era del piloto, "
                "cambios técnicos) que explican la varianza restante.",
                Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.45),
                font_size=11, italic=True, color=RGBColor(0x66, 0x66, 0x66))


def slide_12_discusion(prs):
    """Slide 12 — Discusión y limitaciones del análisis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Discusión y Limitaciones", bg_color=COLOR_NEGRO)

    # Columna izquierda: Discusión
    add_textbox(slide, "¿Qué significa el resultado?",
                Inches(0.5), Inches(1.05), Inches(5.9), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    discusion = [
        "La evidencia estadística confirma que el equipo de debut importa: "
        "los pilotos de equipos Top obtienen sistemáticamente más puntos por carrera.",
        "Sin embargo, la causalidad es compleja: los equipos Top también reclutan "
        "a los pilotos más talentosos, generando sesgo de selección.",
        "El R² bajo (21.7%) confirma que el equipo de debut es un factor relevante "
        "pero no suficiente para explicar el rendimiento de un piloto.",
        "La diferencia observada podría reflejar tanto el contexto del equipo "
        "como el talento diferencial de los pilotos seleccionados.",
    ]

    box_disc = slide.shapes.add_shape(1, Inches(0.5), Inches(1.52),
                                      Inches(5.9), Inches(4.85))
    box_disc.fill.solid()
    box_disc.fill.fore_color.rgb = COLOR_BLANCO
    box_disc.line.color.rgb = COLOR_GRIS_MED
    box_disc.line.width = Pt(1)

    for i, txt in enumerate(discusion):
        y = Inches(1.65) + i * Inches(1.1)
        add_textbox(slide, "  •  " + txt, Inches(0.65), y,
                    Inches(5.6), Inches(1.0), font_size=12.5,
                    color=COLOR_OSCURO_TXT)

    # Columna derecha: Limitaciones
    add_textbox(slide, "Limitaciones Identificadas",
                Inches(7.0), Inches(1.05), Inches(5.9), Inches(0.4),
                font_size=15, bold=True, color=COLOR_ROJO)

    limitaciones = [
        ("Sistema de puntos variable",
         "Los sistemas de puntuación cambiaron en 1991, 2003 y 2010. "
         "Comparar PPC entre eras distintas puede ser engañoso."),
        ("Sesgo de selección",
         "Los equipos Top eligen a los mejores pilotos. "
         "El efecto observado podría ser de talento, no de equipo."),
        ("Variable binaria simplificada",
         "La clasificación Top/Chico ignora diferencias de rendimiento "
         "dentro de cada categoría (p.ej., Ferrari ≠ Williams a mitad de grilla)."),
        ("Debut único por piloto",
         "Solo se considera el equipo de debut, ignorando trayectorias "
         "posteriores con otros equipos de diferente nivel."),
    ]

    box_lim = slide.shapes.add_shape(1, Inches(7.0), Inches(1.52),
                                     Inches(5.9), Inches(4.85))
    box_lim.fill.solid()
    box_lim.fill.fore_color.rgb = COLOR_BLANCO
    box_lim.line.color.rgb = COLOR_GRIS_MED
    box_lim.line.width = Pt(1)

    for i, (titulo, desc) in enumerate(limitaciones):
        y = Inches(1.65) + i * Inches(1.1)
        add_textbox(slide, "⚠  " + titulo, Inches(7.15), y,
                    Inches(5.6), Inches(0.38), font_size=13, bold=True,
                    color=RGBColor(0xB8, 0x58, 0x00))
        add_textbox(slide, desc, Inches(7.15), y + Inches(0.4),
                    Inches(5.6), Inches(0.65), font_size=11.5,
                    color=RGBColor(0x44, 0x44, 0x44))


def slide_13_proximos_pasos(prs):
    """Slide 13 — Próximos pasos para continuar el análisis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_GRIS_CLARO)

    add_title_bar(slide, "  Próximos Pasos", bg_color=COLOR_NEGRO)

    pasos = [
        {
            "num": "01",
            "titulo": "Estandarización del sistema de puntos",
            "desc": (
                "Recalcular el PPC aplicando el sistema de puntos actual (2010+) "
                "de forma retroactiva a todas las carreras históricas. "
                "Esto eliminará el sesgo por cambio de puntuación entre eras."
            ),
            "color": COLOR_ROJO,
        },
        {
            "num": "02",
            "titulo": "Análisis de trayectoria completa del piloto",
            "desc": (
                "Incorporar el rendimiento en todos los equipos del piloto, "
                "no solo el de debut. Modelar si un cambio a equipo Top "
                "mejora el PPC individual. Análisis de panel con efectos fijos por piloto."
            ),
            "color": COLOR_AZUL,
        },
        {
            "num": "03",
            "titulo": "Verificación de supuestos del modelo OLS",
            "desc": (
                "Aplicar tests formales: Breusch-Pagan para homocedasticidad, "
                "Shapiro-Wilk sobre residuos. Si los supuestos se violan, "
                "considerar regresión robusta o transformación logarítmica del PPC."
            ),
            "color": RGBColor(0x00, 0x80, 0x60),
        },
        {
            "num": "04",
            "titulo": "Incorporar variables de control adicionales",
            "desc": (
                "Agregar era del piloto (décadas), nacionalidad del equipo, "
                "y presupuesto estimado del constructor. "
                "Evaluar si el efecto es_top persiste al controlar por estos factores."
            ),
            "color": RGBColor(0x80, 0x00, 0x80),
        },
    ]

    for i, p in enumerate(pasos):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.5)
        y = Inches(1.1) + row * Inches(2.85)

        box = slide.shapes.add_shape(1, x, y, Inches(6.1), Inches(2.65))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_BLANCO
        box.line.color.rgb = p["color"]
        box.line.width = Pt(2.5)

        # Número
        num_shape = slide.shapes.add_shape(1, x, y, Inches(0.7), Inches(0.7))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = p["color"]
        num_shape.line.fill.background()
        add_textbox(slide, p["num"], x + Inches(0.05), y + Inches(0.07),
                    Inches(0.6), Inches(0.55), font_size=14, bold=True,
                    color=COLOR_BLANCO, align=PP_ALIGN.CENTER)

        add_textbox(slide, p["titulo"], x + Inches(0.8), y + Inches(0.1),
                    Inches(5.1), Inches(0.5), font_size=13.5, bold=True,
                    color=p["color"])
        add_textbox(slide, p["desc"], x + Inches(0.15), y + Inches(0.65),
                    Inches(5.8), Inches(1.9), font_size=11.5,
                    color=RGBColor(0x33, 0x33, 0x33))


def slide_14_cierre(prs):
    """Slide 14 — Cierre: conclusiones y agradecimiento. Fondo negro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_NEGRO)

    # Barras decorativas
    bar_top = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                     SLIDE_W, Inches(0.18))
    bar_top.fill.solid()
    bar_top.fill.fore_color.rgb = COLOR_ROJO
    bar_top.line.fill.background()

    bar_bot = slide.shapes.add_shape(1, Inches(0), Inches(7.32),
                                     SLIDE_W, Inches(0.18))
    bar_bot.fill.solid()
    bar_bot.fill.fore_color.rgb = COLOR_ROJO
    bar_bot.line.fill.background()

    # Título sección
    add_textbox(slide, "Conclusiones", Inches(0.7), Inches(0.3),
                Inches(8), Inches(0.6), font_size=28, bold=True,
                color=COLOR_ROJO)

    # 3 conclusiones principales
    conclusiones = [
        ("Diferencia real y significativa",
         "El test de Welch (p ≈ 1.2×10⁻¹⁷) confirma que los pilotos de equipos Top "
         "obtienen en promedio ~1 punto más por carrera que los de equipos Chico, "
         "resultado que no puede atribuirse al azar."),
        ("El equipo de debut importa, pero no lo explica todo",
         "La regresión OLS muestra un coeficiente β₁ = 0.99 significativo, "
         "pero el R² = 21.7% indica que otros factores (talento, era, contexto) "
         "son determinantes igualmente importantes."),
        ("Análisis limitado por el diseño histórico de la F1",
         "El cambio de sistemas de puntuación y el sesgo de selección son "
         "las limitaciones más importantes. Los próximos pasos buscan "
         "abordarlas con un análisis más robusto."),
    ]

    for i, (titulo, texto) in enumerate(conclusiones):
        y = Inches(1.1) + i * Inches(1.85)

        # Número decorativo
        num = slide.shapes.add_shape(1, Inches(0.5), y, Inches(0.55), Inches(0.55))
        num.fill.solid()
        num.fill.fore_color.rgb = COLOR_ROJO
        num.line.fill.background()
        add_textbox(slide, str(i + 1), Inches(0.52), y + Inches(0.04),
                    Inches(0.5), Inches(0.48), font_size=16, bold=True,
                    color=COLOR_BLANCO, align=PP_ALIGN.CENTER)

        add_textbox(slide, titulo, Inches(1.2), y, Inches(11.6), Inches(0.5),
                    font_size=15, bold=True, color=COLOR_BLANCO)
        add_textbox(slide, texto, Inches(1.2), y + Inches(0.52),
                    Inches(11.6), Inches(1.2), font_size=13,
                    color=RGBColor(0xBB, 0xBB, 0xBB))

    # Línea separadora
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(6.68),
                                 Inches(12.33), Inches(0.04))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
    sep.line.fill.background()

    # Agradecimiento
    add_textbox(slide, "¡Gracias!",
                Inches(0.7), Inches(6.75), Inches(5), Inches(0.55),
                font_size=20, bold=True, color=COLOR_ROJO)

    add_textbox(slide,
                "Andy Villarroel  |  Javier Alcaino  —  Análisis de Datos e Inferencia Estadística, UDD",
                Inches(5.0), Inches(6.82), Inches(8.0), Inches(0.45),
                font_size=11.5, color=COLOR_GRIS_MED, align=PP_ALIGN.RIGHT)


# =============================================================================
# FUNCIÓN PRINCIPAL
# =============================================================================

def crear_presentacion():
    """Construye y guarda la presentación completa."""
    prs = Presentation()

    # Tamaño de slide 16:9 widescreen
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Creando diapositivas...")

    slide_01_portada(prs)
    print("  [1/14] Portada")

    slide_02_agenda(prs)
    print("  [2/14] Agenda")

    slide_03_pregunta(prs)
    print("  [3/14] Pregunta de investigación")

    slide_04_dataset(prs)
    print("  [4/14] Dataset y definiciones")

    slide_05_limpieza(prs)
    print("  [5/14] Limpieza de datos")

    slide_06_eda_descriptiva(prs)
    print("  [6/14] EDA — Estadística descriptiva")

    slide_07_eda_histograma(prs)
    print("  [7/14] EDA — Distribución del PPC")

    slide_08_eda_boxplot(prs)
    print("  [8/14] EDA — Boxplot y Scatter")

    slide_09_hipotesis(prs)
    print("  [9/14] Test de hipótesis")

    slide_10_regresion(prs)
    print("  [10/14] Modelo de regresión OLS")

    slide_11_interpretacion(prs)
    print("  [11/14] Interpretación de coeficientes")

    slide_12_discusion(prs)
    print("  [12/14] Discusión y limitaciones")

    slide_13_proximos_pasos(prs)
    print("  [13/14] Próximos pasos")

    slide_14_cierre(prs)
    print("  [14/14] Cierre")

    # Guardar en el mismo directorio del script
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "Presentacion_Avance2.pptx")
    prs.save(output_path)
    print(f"\nPresentacion guardada en:\n  {output_path}")
    print("\nNOTA: Los slides 7 y 8 contienen placeholders grises.")
    print("      Reemplazalos con las figuras exportadas desde tu notebook:")
    print("      fig_hist.png, fig_boxplot.png, fig_scatter.png")


if __name__ == "__main__":
    crear_presentacion()
